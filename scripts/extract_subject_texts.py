#!/usr/bin/env python3
"""Extract text snippets for every asset stored under ``subjects/``.

The script walks the source ``subjects`` directory, extracts textual content
where possible, and saves the normalised output into
``src/data/subjectExtracts``. Each generated ``.txt`` file mirrors the
relative path of the original asset so downstream tooling can provide in-app
previews or search across the raw course materials.
"""
from __future__ import annotations

import argparse
import base64
import binascii
import json
import subprocess
import textwrap
import re
import sys
import zipfile
from collections import defaultdict
from dataclasses import asdict, dataclass
from io import StringIO
from pathlib import Path
import shutil
from typing import Callable, Sequence

try:  # pragma: no cover - optional dependency may be missing in CI environments
    from docx import Document  # type: ignore
except ImportError as docx_import_error:  # pragma: no cover - fallback when dependency absent
    Document = None  # type: ignore[assignment]
    DOCX_IMPORT_ERROR = docx_import_error
else:
    DOCX_IMPORT_ERROR = None
try:  # pragma: no cover - optional dependency may be missing in CI environments
    from openpyxl import load_workbook  # type: ignore
except ImportError as openpyxl_import_error:  # pragma: no cover - fallback when dependency absent
    load_workbook = None  # type: ignore[assignment]
    OPENPYXL_IMPORT_ERROR = openpyxl_import_error
else:
    OPENPYXL_IMPORT_ERROR = None

try:  # pragma: no cover - optional dependency may be missing in CI environments
    from pptx import Presentation  # type: ignore
except ImportError as pptx_import_error:  # pragma: no cover - fallback when dependency absent
    Presentation = None  # type: ignore[assignment]
    PPTX_IMPORT_ERROR = pptx_import_error
else:
    PPTX_IMPORT_ERROR = None
try:  # pragma: no cover - optional dependency may be missing in CI environments
    from pypdf import PdfReader
except ImportError as pypdf_import_error:  # pragma: no cover - fallback when dependency absent
    PdfReader = None  # type: ignore[assignment]
    PYPDF_IMPORT_ERROR = pypdf_import_error
else:
    PYPDF_IMPORT_ERROR = None
PYPDF_AUTOINSTALL_ATTEMPTED = False
try:  # pragma: no cover - optional dependency may be missing in CI environments
    import xlrd  # type: ignore
except ImportError as xlrd_import_error:  # pragma: no cover - fallback when dependency absent
    xlrd = None  # type: ignore[assignment]
    XLRD_IMPORT_ERROR = xlrd_import_error
else:
    XLRD_IMPORT_ERROR = None

ROOT = Path(__file__).resolve().parents[1]
SUBJECTS_DIR = ROOT / "subjects"
OUTPUT_DIR = ROOT / "src" / "data" / "subjectExtracts"
PUBLIC_ASSETS_DIR = ROOT / "public" / "subject-assets"


@dataclass
class ExtractionResult:
    text: str
    notes: list[str]


@dataclass
class ImageMetadata:
    path: str
    page: int
    index: int
    width: int | None
    height: int | None
    color_space: str | None


def _ensure_pypdf_available() -> bool:
    """Attempt to make the ``pypdf`` dependency available on demand."""

    global PdfReader, PYPDF_IMPORT_ERROR, PYPDF_AUTOINSTALL_ATTEMPTED  # type: ignore[assignment]

    if PdfReader is not None:
        return True

    if PYPDF_AUTOINSTALL_ATTEMPTED:
        return False

    PYPDF_AUTOINSTALL_ATTEMPTED = True

    install_cmd = [sys.executable, "-m", "pip", "install", "pypdf"]
    print(
        "[extract-subject-texts] Attempting to automatically install optional dependency 'pypdf'.",
        file=sys.stderr,
    )

    try:
        subprocess.check_call(
            install_cmd,
            stdout=subprocess.DEVNULL,
            stderr=subprocess.STDOUT,
        )
    except Exception as error:  # pragma: no cover - network/tools may be unavailable
        PYPDF_IMPORT_ERROR = RuntimeError(
            "Automatic installation of pypdf failed; install it manually to enable PDF extraction."
            f" (command: {' '.join(install_cmd)}; error: {error})"
        )
        return False

    try:  # pragma: no cover - depends on external installation step
        from pypdf import PdfReader as imported_reader  # type: ignore
    except Exception as import_error:  # pragma: no cover - defensive fallback
        PYPDF_IMPORT_ERROR = import_error
        return False

    PdfReader = imported_reader  # type: ignore[assignment]
    PYPDF_IMPORT_ERROR = None
    return True


def _normalise_whitespace(text: str) -> str:
    """Collapse noisy whitespace while preserving intentional spacing."""
    cleaned = re.sub(r"\r\n?", "\n", text)
    cleaned = re.sub(r"[\t\u00a0]+", " ", cleaned)
    cleaned = re.sub(r"\u200b", "", cleaned)
    cleaned = re.sub(r"\n{3,}", "\n\n", cleaned)
    cleaned = re.sub(r"[ \t]+\n", "\n", cleaned)
    return cleaned.strip()


def _collect_page_image_references(pdf_path: Path) -> dict[int, list[str]]:
    """Return ordered image references for each page of ``pdf_path``."""

    images_dir = pdf_path.with_name(f"{pdf_path.stem}-images")
    if not images_dir.exists():
        return {}

    pattern = re.compile(
        rf"^{re.escape(pdf_path.stem)}_page_(?P<page>\d{{3}})(?:_img_(?P<img>\d{{3}}))?\.png$",
        re.IGNORECASE,
    )
    grouped: dict[int, list[tuple[tuple[int, int], Path]]] = defaultdict(list)

    for image_path in sorted(images_dir.iterdir()):
        if not image_path.is_file() or image_path.suffix.lower() != ".png":
            continue
        match = pattern.match(image_path.name)
        if not match:
            continue
        page_number = int(match.group("page"))
        image_index = match.group("img")
        order_key = (0, 0) if image_index is None else (1, int(image_index))
        grouped[page_number].append((order_key, image_path))

    references: dict[int, list[str]] = {}
    for page_number, entries in grouped.items():
        ordered_paths = [path for _, path in sorted(entries, key=lambda item: item[0])]
        references[page_number] = []
        for image_path in ordered_paths:
            try:
                relative_path = image_path.relative_to(ROOT)
            except ValueError:
                relative_path = image_path
            references[page_number].append(relative_path.as_posix())

    return references


def _resolve_public_asset_dir(pdf_path: Path) -> Path:
    try:
        relative = pdf_path.relative_to(SUBJECTS_DIR)
        target = relative.with_suffix("")
    except ValueError:  # pragma: no cover - fallback for unexpected locations
        target = Path(pdf_path.stem)
    return PUBLIC_ASSETS_DIR / target


def _extract_images_to_public_assets(
    pdf_path: Path,
) -> tuple[dict[int, list[str]], list[ImageMetadata]]:
    try:
        from pdf_image_extractor import extract_images as extract_pdf_images
    except (ImportError, SystemExit):  # pragma: no cover - dependency guard
        return {}, []

    target_dir = _resolve_public_asset_dir(pdf_path)
    if target_dir.exists():
        shutil.rmtree(target_dir)
    target_dir.mkdir(parents=True, exist_ok=True)

    try:
        raw_metadata = extract_pdf_images(pdf_path, target_dir)
    except Exception:  # pragma: no cover - extraction robustness
        return {}, []

    page_references: dict[int, list[str]] = defaultdict(list)
    metadata: list[ImageMetadata] = []
    page_counters: dict[int, int] = defaultdict(int)

    for page_number, filename in sorted(raw_metadata, key=lambda item: (item[0], item[1])):
        image_path = target_dir / filename
        try:
            relative_path = image_path.relative_to(ROOT).as_posix()
        except ValueError:  # pragma: no cover - unexpected outside repo
            relative_path = image_path.as_posix()

        page_counters[page_number] += 1
        page_references[page_number].append(relative_path)
        metadata.append(
            ImageMetadata(
                path=relative_path,
                page=page_number,
                index=page_counters[page_number],
                width=None,
                height=None,
                color_space=None,
            )
        )

    return page_references, metadata


def _format_markdown_image_path(path: str) -> str:
    normalised = path.replace("\\", "/")
    if normalised.startswith("public/"):
        return f"../../{normalised}"
    if normalised.startswith("./"):
        return normalised[2:]
    return normalised


def _initialise_image_output_dir(base_dir: Path, pdf_path: Path) -> Path:
    target_dir = base_dir / pdf_path.stem
    if target_dir.exists():
        shutil.rmtree(target_dir)
    target_dir.mkdir(parents=True, exist_ok=True)
    return target_dir


def _copy_fallback_page_images(
    pdf_path: Path, page_number: int, target_dir: Path
) -> tuple[list[str], list[ImageMetadata]]:
    fallback_dir = pdf_path.parent / f"{pdf_path.stem}-images"
    if not fallback_dir.exists():
        return [], []

    pattern = re.compile(
        rf"^page_{page_number:03d}_img_(?P<index>\d{{3}})\.(?P<ext>png|jpe?g|gif)(?P<encoding>\.base64|\.b64)?$",
        re.IGNORECASE,
    )
    references: list[str] = []
    metadata: list[ImageMetadata] = []

    for source in sorted(fallback_dir.iterdir()):
        if not source.is_file():
            continue

        match = pattern.match(source.name)
        if not match:
            continue

        image_index = int(match.group("index"))
        extension = match.group("ext")
        encoded_suffix = match.group("encoding") or ""
        destination = target_dir / f"page_{page_number:03d}_img_{image_index:03d}.{extension}"
        destination.parent.mkdir(parents=True, exist_ok=True)

        try:
            if encoded_suffix:
                raw_data = source.read_text(encoding="utf-8").strip()
                binary = base64.b64decode(raw_data, validate=True)
                with destination.open("wb") as handle:
                    handle.write(binary)
            else:
                shutil.copy2(source, destination)
        except (OSError, binascii.Error, ValueError):
            continue

        try:
            relative_path = destination.relative_to(ROOT)
        except ValueError:  # pragma: no cover - unexpected outside repo
            relative_path = destination

        references.append(relative_path.as_posix())
        metadata.append(
            ImageMetadata(
                path=relative_path.as_posix(),
                page=page_number,
                index=image_index,
                width=None,
                height=None,
                color_space=None,
            )
        )

    return references, metadata


def _store_page_images(
    page, page_number: int, target_dir: Path, pdf_path: Path
) -> tuple[list[str], list[ImageMetadata]]:
    references: list[str] = []
    metadata: list[ImageMetadata] = []
    try:
        images_iterable = list(getattr(page, "images", []))
    except Exception:  # pragma: no cover - defensive fallback
        images_iterable = []

    if not images_iterable:
        return _copy_fallback_page_images(pdf_path, page_number, target_dir)

    for image_index, image in enumerate(images_iterable, start=1):
        extension = getattr(image, "ext", None) or getattr(image, "extension", None) or "png"
        extension = extension.lstrip(".") or "png"
        filename = f"page_{page_number:03d}_img_{image_index:03d}.{extension}"
        output_path = target_dir / filename
        try:
            with output_path.open("wb") as handle:
                handle.write(image.data)
        except Exception:  # pragma: no cover - avoid breaking extraction on failure
            continue

        try:
            relative_path = output_path.relative_to(ROOT)
        except ValueError:  # pragma: no cover - unexpected outside repo
            relative_path = output_path

        references.append(relative_path.as_posix())
        metadata.append(
            ImageMetadata(
                path=relative_path.as_posix(),
                page=page_number,
                index=image_index,
                width=getattr(image, "width", None),
                height=getattr(image, "height", None),
                color_space=getattr(image, "color_space", None),
            )
        )

    return references, metadata


def _extract_pdf_with_optional_images(
    path: Path, *, image_output_dir: Path | None = None
) -> tuple[ExtractionResult, list[ImageMetadata]]:
    if PdfReader is None and not _ensure_pypdf_available():
        notes = [
            "pypdf is not installed; PDF content was not extracted.",
        ]
        if PYPDF_IMPORT_ERROR is not None:
            notes.append(f"Import error: {PYPDF_IMPORT_ERROR}")
        return (
            ExtractionResult("[PDF extraction requires pypdf to be installed]", notes),
            [],
        )

    if PdfReader is None:  # pragma: no cover - defensive guard
        return (
            ExtractionResult(
                "[PDF extraction requires pypdf to be installed]",
                [
                    "pypdf dependency could not be initialised despite auto-install attempt.",
                    f"Import error: {PYPDF_IMPORT_ERROR}",
                ],
            ),
            [],
        )

    reader = PdfReader(path)
    pieces: list[str] = []
    metadata: list[ImageMetadata] = []

    if image_output_dir is None:
        page_images, collected_metadata = _extract_images_to_public_assets(path)
        metadata.extend(collected_metadata)
        target_dir: Path | None = None
    else:
        target_dir = _initialise_image_output_dir(image_output_dir, path)
        page_images = {}

    for index, page in enumerate(reader.pages, start=1):
        text = (page.extract_text() or "").strip()
        if target_dir is not None:
            images, page_metadata = _store_page_images(page, index, target_dir, path)
            metadata.extend(page_metadata)
        else:
            images = page_images.get(index, [])
        if not text and not images:
            continue

        page_lines = [f"### Page {index}"]
        if text:
            page_lines.append(text)
        for figure_index, image_path in enumerate(images, start=1):
            page_lines.append(
                f"![Page {index}, Figure {figure_index}]({_format_markdown_image_path(image_path)})"
            )

        pieces.append("\n".join(page_lines))

    if not pieces:
        return (
            ExtractionResult(
                "[No text content extracted]",
                ["PDF parser returned no text; file may be scanned images."],
            ),
            metadata,
        )

    return ExtractionResult("\n\n".join(pieces), []), metadata


def extract_pdf(path: Path, *, image_output_dir: Path | None = None) -> ExtractionResult:
    result, _ = _extract_pdf_with_optional_images(path, image_output_dir=image_output_dir)
    return result


def extract_ipynb(path: Path) -> ExtractionResult:
    data = json.loads(path.read_text(encoding="utf-8"))
    parts: list[str] = []
    for index, cell in enumerate(data.get("cells", []), start=1):
        cell_type = cell.get("cell_type", "unknown")
        header = f"### Cell {index} · {cell_type.title()}"
        source = "".join(cell.get("source", []))
        if cell_type == "code":
            source = source.rstrip()
            parts.append(f"{header}\n```python\n{source}\n```")
        else:
            parts.append(f"{header}\n{source}")
    if not parts:
        return ExtractionResult("[Notebook contains no cells]", [])
    return ExtractionResult("\n\n".join(parts), [])


def extract_text_file(path: Path) -> ExtractionResult:
    return ExtractionResult(path.read_text(encoding="utf-8", errors="replace"), [])


def extract_url(path: Path) -> ExtractionResult:
    content = path.read_text(encoding="utf-8", errors="replace")
    return ExtractionResult(content.strip(), [])


def extract_presentation(path: Path) -> ExtractionResult:
    if Presentation is None:
        notes = [
            "python-pptx is not installed; PPTX/PPSX content was not extracted.",
        ]
        if PPTX_IMPORT_ERROR is not None:
            notes.append(f"Import error: {PPTX_IMPORT_ERROR}")
        return ExtractionResult("[Presentation extraction requires python-pptx to be installed]", notes)

    presentation = Presentation(path)
    pieces: list[str] = []
    for slide_number, slide in enumerate(presentation.slides, start=1):
        slide_parts: list[str] = []
        for shape in slide.shapes:
            if getattr(shape, "has_text_frame", False):
                text = shape.text.strip()
                if text:
                    slide_parts.append(text)
            if getattr(shape, "has_table", False):
                rows = []
                for row in shape.table.rows:
                    cells = [cell.text.strip() for cell in row.cells]
                    if any(cells):
                        rows.append(" | ".join(cells))
                if rows:
                    slide_parts.append("Table:\n" + "\n".join(rows))
        if slide_parts:
            pieces.append(f"### Slide {slide_number}\n" + "\n\n".join(slide_parts))
    if not pieces:
        return ExtractionResult("[No text content extracted from presentation]", [])
    return ExtractionResult("\n\n".join(pieces), [])


def extract_excel_xlsx(path: Path) -> ExtractionResult:
    if load_workbook is None:
        notes = [
            "openpyxl is not installed; XLSX content was not extracted.",
        ]
        if OPENPYXL_IMPORT_ERROR is not None:
            notes.append(f"Import error: {OPENPYXL_IMPORT_ERROR}")
        return ExtractionResult("[XLSX extraction requires openpyxl to be installed]", notes)

    workbook = load_workbook(path, data_only=True, read_only=True)
    buffer = StringIO()
    for sheet in workbook.worksheets:
        buffer.write(f"### Sheet: {sheet.title}\n")
        for row in sheet.iter_rows(values_only=True):
            values = ["" if value is None else str(value) for value in row]
            if any(value.strip() for value in values):
                buffer.write("\t".join(values))
                buffer.write("\n")
        buffer.write("\n")
    text = buffer.getvalue().strip()
    if not text:
        return ExtractionResult("[No cell content extracted from workbook]", [])
    return ExtractionResult(text, [])


def extract_excel_xls(path: Path) -> ExtractionResult:
    if xlrd is None:
        notes = [
            "xlrd is not installed; XLS content was not extracted.",
        ]
        if XLRD_IMPORT_ERROR is not None:
            notes.append(f"Import error: {XLRD_IMPORT_ERROR}")
        return ExtractionResult("[XLS extraction requires xlrd to be installed]", notes)

    workbook = xlrd.open_workbook(path)
    buffer = StringIO()
    for sheet in workbook.sheets():
        buffer.write(f"### Sheet: {sheet.name}\n")
        for row_index in range(sheet.nrows):
            row_values = [sheet.cell_value(row_index, col_index) for col_index in range(sheet.ncols)]
            formatted = ["" if value in ("", None) else str(value) for value in row_values]
            if any(value.strip() for value in formatted):
                buffer.write("\t".join(formatted))
                buffer.write("\n")
        buffer.write("\n")
    text = buffer.getvalue().strip()
    if not text:
        return ExtractionResult("[No cell content extracted from workbook]", [])
    return ExtractionResult(text, [])


def extract_docx(path: Path) -> ExtractionResult:
    if Document is None:
        notes = [
            "python-docx is not installed; DOCX content was not extracted.",
        ]
        if DOCX_IMPORT_ERROR is not None:
            notes.append(f"Import error: {DOCX_IMPORT_ERROR}")
        return ExtractionResult("[DOCX extraction requires python-docx to be installed]", notes)

    document = Document(path)
    parts = [para.text for para in document.paragraphs if para.text.strip()]
    for table in document.tables:
        for row in table.rows:
            cells = [cell.text.strip() for cell in row.cells]
            if any(cells):
                parts.append(" | ".join(cells))
    if not parts:
        return ExtractionResult("[No text extracted from document]", [])
    return ExtractionResult("\n".join(parts), [])


def extract_archimate(path: Path) -> ExtractionResult:
    notes: list[str] = []
    pieces: list[str] = []
    try:
        with zipfile.ZipFile(path) as archive:
            for name in sorted(archive.namelist()):
                if not name.lower().endswith((".xml", ".txt")):
                    continue
                data = archive.read(name).decode("utf-8", errors="replace")
                pieces.append(f"### {name}\n{data.strip()}")
    except zipfile.BadZipFile:
        notes.append("ArchiMate file is not a ZIP archive; storing binary notice instead.")
    if not pieces:
        return ExtractionResult("[No textual XML resources extracted]", notes)
    return ExtractionResult("\n\n".join(pieces), notes)


def extract_png(path: Path) -> ExtractionResult:
    return ExtractionResult("[Binary image file – no textual content extracted]", [])


def extract_generic(path: Path) -> ExtractionResult:
    return ExtractionResult(f"[Unsupported file type: {path.suffix}]", [])


EXTRACTORS: dict[str, Callable[[Path], ExtractionResult]] = {
    ".pdf": extract_pdf,
    ".ipynb": extract_ipynb,
    ".txt": extract_text_file,
    ".sql": extract_text_file,
    ".url": extract_url,
    ".pptx": extract_presentation,
    ".ppsx": extract_presentation,
    ".xlsx": extract_excel_xlsx,
    ".xls": extract_excel_xls,
    ".docx": extract_docx,
    ".archimate": extract_archimate,
    ".png": extract_png,
}


def extract_file(path: Path) -> ExtractionResult:
    extractor = EXTRACTORS.get(path.suffix.lower())
    if extractor is None:
        return extract_generic(path)
    try:
        return extractor(path)
    except Exception as error:  # noqa: BLE001 - pipeline must be resilient
        return ExtractionResult(
            text=f"[Failed to extract content: {error}]",
            notes=[f"Extraction error for {path.name}: {error}"],
        )


def build_header(relative_path: Path, notes: list[str]) -> str:
    header_lines = [
        "# Extracted content",
        f"Source: subjects/{relative_path.as_posix()}",
    ]
    if notes:
        header_lines.append("Notes:")
        header_lines.extend(f"- {note}" for note in notes)
    return "\n".join(header_lines) + "\n\n"


def _write_support_modules() -> None:
    glob_modules_path = OUTPUT_DIR / "globModules.ts"
    glob_modules_path.write_text(
        "\n".join(
            [
                "const subjectExtractModules = import.meta.glob('./**/*.txt', {",
                "  eager: true,",
                "  import: 'default',",
                "  query: '?raw',",
                "}) as Record<string, string>;",
                "",
                "export default subjectExtractModules;",
            ]
        )
        + "\n",
        encoding="utf-8",
    )

    index_module_path = OUTPUT_DIR / "index.ts"
    index_module_path.write_text(
        "\n".join(
            [
                "import subjectExtractModules from './globModules';",
                "",
                "const loadModulesWithFs = (): Record<string, string> => {",
                "  try {",
                "    // eslint-disable-next-line @typescript-eslint/no-var-requires",
                "    const fs = require('fs') as typeof import('fs');",
                "    // eslint-disable-next-line @typescript-eslint/no-var-requires",
                "    const path = require('path') as typeof import('path');",
                "    const rootDir = __dirname;",
                "    const result: Record<string, string> = {};",
                "",
                "    const visit = (dir: string) => {",
                "      const entries = fs.readdirSync(dir, { withFileTypes: true });",
                "      for (const entry of entries) {",
                "        const entryPath = path.join(dir, entry.name);",
                "        if (entry.isDirectory()) {",
                "          visit(entryPath);",
                "        } else if (entry.isFile() && entry.name.endsWith('.txt')) {",
                r"          const relativePath = `./${path.relative(rootDir, entryPath).replace(/\\/g, '/')}`;",
                "          result[relativePath] = fs.readFileSync(entryPath, 'utf8');",
                "        }",
                "      }",
                "    };",
                "",
                "    visit(rootDir);",
                "    return result;",
                "  } catch (error) {",
                "    console.warn('[subjectExtracts] Unable to load extracted text via fs:', error);",
                "    return {};",
                "  }",
                "};",
                "",
                "const modules = (() => {",
                "  if (subjectExtractModules && Object.keys(subjectExtractModules).length > 0) {",
                "    return subjectExtractModules;",
                "  }",
                "  return loadModulesWithFs();",
                "})();",
                "",
                "type ExtractedSubjectText = {",
                "  /** Path of the source asset inside the `subjects/` tree. */",
                "  source: string;",
                "  /** Normalised text extracted from the source asset, excluding the metadata header. */",
                "  text: string;",
                "  /** Optional extraction notes declared in the file header. */",
                "  notes?: string[];",
                "  /** Absolute module id used by Vite (useful for debugging). */",
                "  moduleId: string;",
                "};",
                "",
                r"const headerRegex = /^# Extracted content\nSource: (?<source>subjects\/[\s\S]+?)\n(?:Notes:\n(?<notes>(?:- .+\n)+))?/;",
                "",
                "const map = new Map<string, ExtractedSubjectText>();",
                "",
                "for (const [moduleId, raw] of Object.entries(modules)) {",
                r"  const rawText = raw.replace(/\r\n/g, '\n');",
                "  const match = rawText.match(headerRegex);",
                "  const source = match?.groups?.source?.trim();",
                "",
                "  if (!source) {",
                "    // Skip files that do not follow the expected header format.",
                "    continue;",
                "  }",
                "",
                "  const notesBlock = match?.groups?.notes;",
                r"  const notes = notesBlock",
                r"    ?.split('\n')",
                r"    .map((line) => line.replace(/^-\s*/, '').trim())",
                "    .filter((line) => line.length > 0);",
                "",
                "  const text = rawText.replace(headerRegex, '').trim();",
                "",
                "  map.set(source.toLowerCase(), {",
                "    source,",
                "    text,",
                "    ...(notes && notes.length > 0 ? { notes } : {}),",
                "    moduleId,",
                "  });",
                "}",
                "",
                "export const subjectExtracts = map;",
                "",
                "export const getSubjectExtract = (sourcePath: string): ExtractedSubjectText | undefined =>",
                "  map.get(sourcePath.toLowerCase());",
                "",
                "export type { ExtractedSubjectText };",
            ]
        )
        + "\n",
        encoding="utf-8",
    )


def _run_bulk_extraction() -> int:
    if not SUBJECTS_DIR.exists():
        print("Subjects directory not found.", file=sys.stderr)
        return 1

    if OUTPUT_DIR.exists():
        shutil.rmtree(OUTPUT_DIR)

    OUTPUT_DIR.mkdir(parents=True, exist_ok=True)
    _write_support_modules()

    written = 0
    for source in sorted(SUBJECTS_DIR.rglob("*")):
        if not source.is_file():
            continue
        relative = source.relative_to(SUBJECTS_DIR)
        output_path = OUTPUT_DIR / relative.with_suffix(".txt")
        output_path.parent.mkdir(parents=True, exist_ok=True)

        result = extract_file(source)
        header = build_header(relative, result.notes)
        content = header + _normalise_whitespace(result.text)
        output_path.write_text(content + "\n", encoding="utf-8")
        written += 1

    print(f"Extracted {written} files into {OUTPUT_DIR.relative_to(ROOT)}")
    return 0


def _extract_single_pdf(pdf_path: Path, images_dir: Path) -> int:
    extraction, metadata = _extract_pdf_with_optional_images(
        pdf_path, image_output_dir=images_dir
    )
    payload = {
        "text": extraction.text,
        "images": [asdict(entry) for entry in metadata],
    }
    print(json.dumps(payload, ensure_ascii=False))
    return 0


def main(argv: Sequence[str] | None = None) -> int:
    parser = argparse.ArgumentParser(description="Extract study subject assets")
    parser.add_argument(
        "--single-pdf",
        type=Path,
        help="Extract a single PDF file and output JSON",
    )
    parser.add_argument(
        "--images-dir",
        type=Path,
        help="Directory where extracted PDF images will be stored",
    )
    parser.add_argument(
        "target",
        nargs="?",
        type=Path,
        help=(
            "Backwards compatible alias for --single-pdf. "
            "Supplying a path positional argument is equivalent to using "
            "--single-pdf."
        ),
    )
    args = parser.parse_args(argv)

    if args.single_pdf is not None and args.target is not None:
        parser.error("Specify either --single-pdf or a positional path, not both.")

    pdf_target = args.single_pdf or args.target

    if pdf_target is not None:
        pdf_path: Path = pdf_target
        if not pdf_path.exists():
            print(f"PDF not found: {pdf_path}", file=sys.stderr)
            return 2
        images_dir = args.images_dir or (SUBJECTS_DIR / "tmp-extracted-images")
        return _extract_single_pdf(pdf_path, images_dir)

    return _run_bulk_extraction()


if __name__ == "__main__":
    raise SystemExit(main())
